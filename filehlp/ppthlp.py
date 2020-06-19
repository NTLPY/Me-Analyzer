"""
    PowerPoint Helper

    This header is to extract the infomation included in ppt files.

    Author:   NTLPY
    Creation: 0x01d645e3149f5f50 (UTC)

"""

# import modules
import pptx as _pptx;
import wcuthlp;

# track shape
def track_shape(shape, words : dict):
    if (type(shape) == _pptx.shapes.group.GroupShape):
        for s in shape.shapes:
            track_shape(s, words);
    else:
        if (shape.has_text_frame):
            wcuthlp.cut_words(shape.text, words);


# extract from pptx
def get_from_pptx(filename : str, words : dict):
    file = _pptx.Presentation(filename);
    # enum all pages(slides)
    for p in file.slides:
        #enum all shapes iterately
        for s in p.shapes:
            track_shape(s, words);
    return;
